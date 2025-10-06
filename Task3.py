import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT_DIR = "phishing_package"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Phishing Awareness Training"
    subtitle.text = "Learn how to spot and stop phishing attacks."

    # Slide 2: What is Phishing?
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "What is Phishing?"
    slide.placeholders[1].text = (
        "Phishing is a type of cyber attack where attackers trick you into revealing personal "
        "information by pretending to be a trusted source."
    )

    # Slide 3: Common Signs
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Common Signs of Phishing"
    slide.placeholders[1].text = (
        "- Suspicious email addresses\n"
        "- Urgent or threatening language\n"
        "- Unexpected attachments or links\n"
        "- Too-good-to-be-true offers"
    )

    # Slide 4: Prevention Tips
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "How to Protect Yourself"
    slide.placeholders[1].text = (
        "- Verify the sender before clicking links\n"
        "- Hover over links to check their real destination\n"
        "- Do not share sensitive info via email\n"
        "- Report suspicious emails to IT/security team"
    )

    # Save with timestamp (safe version)
    PPTX_PATH = os.path.join(OUTPUT_DIR, "Phishing_Awareness.pptx")
    base, ext = os.path.splitext(PPTX_PATH)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_path = f"{base}_{ts}{ext}"

    try:
        prs.save(safe_path)
        print(f"  -> Saved: {safe_path}")
    except PermissionError:
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        alt_dir = os.path.join(desktop, "phishing_package")
        os.makedirs(alt_dir, exist_ok=True)
        alt_path = os.path.join(alt_dir, os.path.basename(safe_path))
        prs.save(alt_path)
        print(f"  -> Saved to alternate path: {alt_path}")

def create_quiz():
    quiz_path = os.path.join(OUTPUT_DIR, "quiz.md")
    with open(quiz_path, "w") as f:
        f.write("# Phishing Awareness Quiz\n\n")
        f.write("**Q1:** What is phishing?\n")
        f.write("a) A cyberattack tricking users into sharing data\n")
        f.write("b) A type of fish\n")
        f.write("c) A safe website\n")
        f.write("\nAnswer: a)\n\n")

        f.write("**Q2:** What should you do if an email looks suspicious?\n")
        f.write("a) Click the link to see what happens\n")
        f.write("b) Delete it or report it\n")
        f.write("c) Reply with personal information\n")
        f.write("\nAnswer: b)\n")
    print(f"  -> Saved: {quiz_path}")

def create_sample_email():
    email_path = os.path.join(OUTPUT_DIR, "sample_phish.txt")
    with open(email_path, "w") as f:
        f.write("From: support@fakebank.com\n")
        f.write("Subject: Urgent: Verify your account\n\n")
        f.write("Dear user,\n")
        f.write("We detected unusual activity on your account. Please login immediately using the link below:\n")
        f.write("http://fakebank-login.com\n\n")
        f.write("Failure to verify will result in account suspension.\n")
    print(f"  -> Saved: {email_path}")

def create_readme():
    readme_path = os.path.join(OUTPUT_DIR, "README.md")
    with open(readme_path, "w") as f:
        f.write("# Phishing Awareness Package\n\n")
        f.write("This package contains:\n")
        f.write("- A training PowerPoint presentation\n")
        f.write("- A phishing awareness quiz\n")
        f.write("- A sample phishing email\n")
        f.write("- This README file\n\n")
        f.write("Use this package to train employees about phishing threats.\n")
    print(f"  -> Saved: {readme_path}")

def main():
    print("Phishing Awareness package generator\n")
    print("Creating PowerPoint presentation...")
    create_presentation()

    print("Creating quiz file...")
    create_quiz()

    print("Creating sample phishing email file...")
    create_sample_email()

    print("Creating README file...")
    create_readme()

    abs_path = os.path.abspath(OUTPUT_DIR)
    print(f"\nAll files created successfully in the folder: {abs_path}")

if __name__ == "__main__":
    main()
